import logging
import os

import numpy as np
import streamlit as st

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def title_prompt(content):
    logger.info("Generating chat title")
    text = str(content)
    instruction = f"""Generate a short, concise title (3–6 words) that summarizes the following message.
    Output only the title as plain text. Do not include quotes, punctuation, commentary, or any additional words!
    Only produce the title itself:
    {clean_text(text)}
    """
    return instruction


def summary_prompt(content):
    logger.info("Generating summary")
    text = str(content)
    instruction = f"""Generate a short, concise summary of the content found in the first 3000 characters of the provided text.
    If the text already contains a summary-like section (e.g., “abstract”, “summary”), extract only that section instead.
    The summary must focus strictly on the subject matter, ideas, findings, or events described in the text.
    Do not focus on the author, their background, or personal details.
    For example, emphasize the concepts, methods, results, or arguments presented, not the person presenting them.
    Output only the summary/extracted text as plain text.
    Only produce the summary/extracted text itself:
    {clean_text(text)}
    """
    return instruction


def build_prompt_template(messages, mode, system_prompt=None, summerize=False):
    prompt = []
    if not summerize:
        if not system_prompt.strip():
            system_prompt = "You are a helpful assistant."
        system_content = [{"type": "text", "text": system_prompt}] if mode == "vlm" else system_prompt
        prompt = [{"role": "system", "content": system_content}]
        for msg in messages[:-1]:
            prompt.append(
                {
                    "role": msg.get("role"),
                    "content": [{"type": "text", "text": msg.get("text")}] if mode == "vlm" else msg.get("text"),
                }
            )
        if mode == "vlm":
            prompt.append({"role": "user", "content": messages[-1].get("vlm_content")})
        else:
            prompt.append({"role": "user", "content": messages[-1].get("content")})
    else:
        for msg in messages:
            prompt.append(
                {
                    "role": msg.get("role"),
                    "content": [{"type": "text", "text": msg.get("text")}] if mode == "vlm" else msg.get("text"),
                }
            )
    return prompt


def summarize_history(messages, mode, model_name, max_window_size=4096):
    if not messages:
        return messages
    template = build_prompt_template(messages=messages, mode=mode, summerize=True)

    from utils.api_client import APIClient

    num_tokens = APIClient.count_tokens(model_name=model_name, text=template)
    logger.info(f"Window: {max_window_size}")
    if num_tokens < max_window_size:
        return messages

    message_placeholder = st.empty()
    message_placeholder.markdown("*Trying to keep up with the whole conversation...*")
    st.toast(
        "The conversation became too long, the model will begin to have a hard time remembering all details.",
        duration="long",
    )
    logger.info("Summarizing old conversation")

    new_messages = []
    query = """
    Summarize the entire previous conversation for memory purposes.

    Instructions:
    - This summary is NOT a dialogue and must NOT be treated as chat turns.
    - Do NOT write responses or continue the conversation.
    - Only describe what the user said and what the assistant said.

    Output format (exactly):

    CONVERSATION SUMMARY
    User:
    <concise summary of everything the user asked or stated>

    Assistant:
    <concise summary of everything the assistant explained or replied>

    End of summary.
    """
    template.append({"role": "user", "content": [{"type": "text", "text": query}] if mode == "vlm" else query})

    user_recall = """The following is a condensed memory of the prior conversation.
    Use it only as background context. Respond normally to the user's next message."""

    new_messages.append(
        {
            "role": "user",
            "text": user_recall,
            "content": user_recall,
            "vlm_content": [{"type": "text", "text": user_recall}],
            "medias": [],
            "docs": [],
        }
    )

    response = APIClient.generate(
        model_name=model_name,
        template=template,
        max_tokens=1024,
    )

    new_messages.append(
        {
            "role": "assistant",
            "vlm_content": [{"type": "text", "text": response}],
            "content": response,
            "text": response,
        }
    )
    message_placeholder.empty()
    return new_messages


def extract_keywords(texts, top_n: int = 5):
    """Extract keywords using TF-IDF - works for any language"""
    from sklearn.feature_extraction.text import TfidfVectorizer

    vectorizer = TfidfVectorizer(
        max_features=120,
        min_df=1,
        max_df=0.6,
        token_pattern=r"\b(?:[A-Z]{2,}|\w{3,})\b",
        norm="l2",
        use_idf=True,
        smooth_idf=True,
        sublinear_tf=True,
        ngram_range=(1, 2),
        analyzer="word",
        lowercase=False,
    )

    try:
        tfidf_matrix = vectorizer.fit_transform(texts)
        feature_names = vectorizer.get_feature_names_out()

        keywords_per_chunk = []
        for row in tfidf_matrix:
            scores = row.toarray()[0]
            top_indices = scores.argsort()[-top_n:][::-1]
            keywords = [feature_names[i] for i in top_indices if scores[i] > 0]
            keywords_per_chunk.append(keywords)
        keywords_per_chunk[-1].append("link")
        return keywords_per_chunk
    except Exception as e:
        print(f"TF-IDF failed: {e}")
        return [[w for w in text.split() if len(w) > 3][:top_n] for text in texts]


def cosine_similarity(a, b):
    """Calculate cosine similarity between two vectors"""
    a = np.array(a).flatten()  # Ensure 1D
    b = np.array(b).flatten()  # Ensure 1D

    dot_product = np.dot(a, b)
    norm_a = np.linalg.norm(a)
    norm_b = np.linalg.norm(b)

    if norm_a == 0 or norm_b == 0:
        return 0.0

    return dot_product / (norm_a * norm_b)


def similarity_matching(collection, embedding_model, all_doc_ids, chunks, threshold: float = 0.90):

    if not all_doc_ids or not chunks:
        return False

    # Embed ALL new chunks and average them
    new_chunk_embeddings = embedding_model.encode(chunks)
    new_avg_embedding = np.mean(new_chunk_embeddings, axis=0)

    logger.info(f"Encoded {len(chunks)} new chunks for comparison")

    for doc_id in all_doc_ids:
        try:
            where_filter = {"doc_id": doc_id}

            # Get pre-computed embeddings from ChromaDB
            all_data = collection.get(where=where_filter, include=["embeddings", "metadatas"])

            if len(all_data["embeddings"]) == 0:
                logger.warning(f"No embeddings found for {doc_id}")
                continue

            # Fix: ChromaDB returns list of embeddings, convert properly
            existing_embeddings = np.array(all_data["embeddings"])  # Shape: (num_chunks, embedding_dim)

            # Check shape
            if existing_embeddings.ndim == 1:
                # Single embedding
                existing_avg_embedding = existing_embeddings
            else:
                # Multiple embeddings - average them
                existing_avg_embedding = np.mean(existing_embeddings, axis=0)

            # Calculate cosine similarity
            similarity = cosine_similarity(new_avg_embedding, existing_avg_embedding)

            if similarity >= threshold:
                logger.info(f"Duplicate detected! {similarity:.2%} match with {doc_id}")
                return True

        except Exception as e:
            logger.error(f"Error checking similarity for {doc_id}: {e}")
            logger.exception(e)  # Show full traceback
            continue

    return False


def hierarchical_retrieval(
    collection, embedding_model, chat_docs, query: str, top_n: int = 5, threshold: float = 0.7, user: str = "user"
):

    import json

    all_data = {}
    file_path = os.path.join(os.getenv("RAG_DIR", os.path.join("outputs", "RAG")), "DOCUMENT_METADATA.json")

    try:
        with open(file_path, "r", encoding="utf-8") as f:
            DOCUMENT_METADATA = json.load(f)
    except FileNotFoundError:
        return {}
    except json.JSONDecodeError:
        raise ValueError(f"Failed to load JSON: {file_path}")

    relevant_doc_ids = {}

    if chat_docs:
        chat_doc_metadata = {doc_id: meta for doc_id, meta in DOCUMENT_METADATA.items() if doc_id in chat_docs}

        query_words = set(query.lower().split())
        query_embedding = embedding_model.encode([query])[0]

        doc_scores = []
        for doc_id, meta in chat_doc_metadata.items():  # Only iterate chat docs
            summary = meta.get("summary", "")

            # Score 1: Word matching
            summary_words = set(summary.lower().split())
            word_score = len(query_words & summary_words) / max(len(query_words), 1)

            # Score 2: Semantic similarity
            summary_embedding = embedding_model.encode([summary])[0]

            semantic_score = cosine_similarity(query_embedding, summary_embedding)

            # Combined score (70% semantic, 30% word match)
            combined_score = 0.3 * word_score + 0.7 * semantic_score
            doc_scores.append((doc_id, combined_score))

        # Sort by score and filter by threshold
        doc_scores.sort(key=lambda x: x[1], reverse=True)
        relevant_doc_ids = {
            doc_id: DOCUMENT_METADATA[doc_id]["total_chunks"] for doc_id, score in doc_scores[:4] if score >= 0.3
        }

    if relevant_doc_ids:
        chunk_ids = [
            f"{doc_id}_chunk_{index}"
            for doc_id, total_chunks in relevant_doc_ids.items()
            for index in range(total_chunks)
        ]
        where_filter = {"chunk_id": {"$in": chunk_ids}}
        all_data = collection.get(where=where_filter, include=["metadatas", "documents"])

    if not all_data or not all_data.get("ids"):
        return {}

    chunk_ids = [meta["chunk_id"] for meta in all_data["metadatas"]]

    results = collection.query(
        query_texts=[query], where={"chunk_id": {"$in": chunk_ids}}, n_results=min(50, len(chunk_ids))
    )

    final_ids = []
    final_docs = []
    final_metas = []
    final_distances = []

    # Collect results that pass threshold
    candidates = []
    for i, distance in enumerate(results["distances"][0]):
        if distance <= threshold:
            candidates.append(
                {
                    "id": results["ids"][0][i],
                    "doc": results["documents"][0][i],
                    "meta": results["metadatas"][0][i],
                    "distance": distance,
                }
            )

    # Sort by distance (lowest first = most similar first)
    candidates.sort(key=lambda x: x["distance"])

    # Extract into separate lists
    for candidate in candidates:
        final_ids.append(candidate["id"])
        final_docs.append(candidate["doc"])
        final_metas.append(candidate["meta"])
        final_distances.append(candidate["distance"])

    final_results = {
        "ids": [final_ids[:top_n]],
        "documents": [final_docs[:top_n]],
        "metadatas": [final_metas[:top_n]],
        "distances": [final_distances[:top_n]],
    }
    return final_results


def media_resize(file, max_width=256, fps=15):
    """Resize image or video and return temp path"""
    import tempfile

    file_extension = os.path.splitext(file.name)[1]
    if file.type.startswith("image"):
        from PIL import Image

        img = Image.open(file).convert("RGB")
        if max(img.size) > max_width:
            img.thumbnail((max_width, max_width), Image.Resampling.LANCZOS)
        with tempfile.NamedTemporaryFile(suffix=file_extension, delete=False) as f:
            img.save(f, format="JPEG", quality=85, optimize=True)
            return f.name, "image"
    elif file.type.startswith("video"):
        import cv2

        # Save uploaded file
        with tempfile.NamedTemporaryFile(suffix=file_extension, delete=False) as f:
            f.write(file.read())
            input_path = f.name

        cap = cv2.VideoCapture(input_path)
        w = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
        h = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
        original_fps = cap.get(cv2.CAP_PROP_FPS)

        target_fps = fps
        frame_interval = int(original_fps / target_fps) if original_fps > target_fps else 1

        if w > max_width:
            new_w, new_h = max_width, int(h * (max_width / w))
        else:
            new_w, new_h = w, h

        output_path = tempfile.NamedTemporaryFile(suffix=".mp4", delete=False).name
        out = cv2.VideoWriter(output_path, cv2.VideoWriter_fourcc(*"mp4v"), target_fps, (new_w, new_h))

        frame_count = 0
        while cap.isOpened():
            ret, frame = cap.read()
            if not ret:
                break

            if frame_count % frame_interval == 0:
                resized = cv2.resize(frame, (new_w, new_h))
                out.write(resized)

            frame_count += 1

        cap.release()
        out.release()
        try:
            os.remove(input_path)
        except FileNotFoundError:
            pass

        return output_path, "video"


def clear_temp(list_of_path):
    for path in list_of_path:
        try:
            os.remove(path)
        except FileNotFoundError:
            pass


def make_docs_query(text_content, documents):
    query = f"""Based on the following sources, answer this question:

            Question: {text_content}

            Sources:
            """

    for i, doc in enumerate(documents):
        query += f"\n--- Source {i} ---\n{doc}.\n"
    query += "\nAnswer based on the sources above:"

    return clean_text(query)


def make_rag_query(text_content, relevant_chunks):

    if not relevant_chunks or not relevant_chunks["ids"][0]:
        return text_content

    context = ""
    sources = []

    for i in range(len(relevant_chunks["ids"][0])):
        chunk_id = relevant_chunks["ids"][0][i]
        document = relevant_chunks["documents"][0][i]
        title = chunk_id.rsplit("_", 4)[0]

        if title not in sources:
            sources.append(title)

        context += f"{document}\n\n"

    source_list = ", ".join(sources[:3])
    if len(sources) > 3:
        source_list += f" and {len(sources) - 3} more"
    print(title)
    full_query = f"""<context>\n
    The user has access to uploaded documents: {source_list}\n
    Relevant excerpts:\n
    {context.strip()}
    </context>
    <user_message>\n
    {text_content}\n
    </user_message>\n
    Respond to the user's message. Use the document context only if it's relevant to what they're asking while citing the sources using their titles. If they're just chatting or the context isn't helpful, respond naturally without forcing references to the documents."""

    return clean_text(full_query)


def read_txt(file) -> str:
    """Read plain text file"""
    return file.read().decode("utf-8")


def read_excel(file) -> str:
    """Read Excel file (XLSX, XLS)"""
    try:
        import pandas as pd

        # Read all sheets
        excel_file = pd.ExcelFile(file)
        text = ""
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file, sheet_name=sheet_name)
            text += f"\n\n=== Sheet: {sheet_name} ===\n"
            text += df.to_string(index=False)
        return text
    except Exception as e:
        st.error(f"Error reading Excel: {e}")
        return ""


def read_csv(file) -> str:
    """Read CSV file"""
    try:
        import pandas as pd

        df = pd.read_csv(file)
        return df.to_string(index=False)
    except Exception as e:
        st.error(f"Error reading CSV: {e}")
        return ""


def read_json(file) -> str:
    """Read JSON file"""
    try:
        import json

        data = json.load(file)
        return json.dumps(data, indent=2)
    except Exception as e:
        st.error(f"Error reading JSON: {e}")
        return ""


def append_hidden_links(text, links):
    """Helper: Only append links that aren't already visible in text"""
    links = [link for link in links if link]
    if not links:
        return text

    text_lower = text.lower()
    hidden = []

    for link in links:
        link_clean = link.replace("https://", "").replace("http://", "").replace("www.", "")
        if link_clean.lower() not in text_lower:
            hidden.append(link)

    if hidden:
        text += "\n\n=== Links ===\n" + "\n".join(hidden)

    return text


def read_pdf(file) -> str:
    """Read PDF"""
    try:
        import fitz

        pdf = fitz.open(stream=file.read(), filetype="pdf")

        text = ""
        links = []

        for page in pdf:
            text += page.get_text("text", sort=True)
            links.extend([link["uri"] for link in page.get_links() if "uri" in link])

        pdf.close()
        return append_hidden_links(text, links)

    except Exception as e:
        return f"Error: {e}"


def read_docx(file) -> str:
    """Read Word document"""
    try:
        import docx

        doc = docx.Document(file)

        text = "\n".join([p.text for p in doc.paragraphs])

        links = []
        for rel in doc.part.rels.values():
            if "hyperlink" in rel.reltype:
                links.append(rel.target_ref)

        return append_hidden_links(text, links)

    except Exception as e:
        return f"Error: {e}"


def read_pptx(file) -> str:
    """Read PowerPoint with hidden links"""
    try:
        from pptx import Presentation

        prs = Presentation(file)

        text = ""
        links = []

        for i, slide in enumerate(prs.slides, 1):
            text += f"\n=== Slide {i} ===\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

                if hasattr(shape, "click_action") and shape.click_action.hyperlink:
                    links.append(shape.click_action.hyperlink.address)

                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.hyperlink and run.hyperlink.address:
                                links.append(run.hyperlink.address)

        return append_hidden_links(text, links)

    except Exception as e:
        return f"Error: {e}"


def read_html(file) -> str:
    """Read HTML"""
    try:
        from bs4 import BeautifulSoup

        html = file.read().decode("utf-8")
        soup = BeautifulSoup(html, "html.parser")

        for script in soup(["script", "style"]):
            script.decompose()

        return soup.get_text(separator="\n", strip=True)

    except Exception as e:
        return f"Error: {e}"


def read_markdown(file) -> str:
    """Read Markdown"""
    try:
        return file.read().decode("utf-8")
    except Exception as e:
        return f"Error: {e}"


def read_rtf(file) -> str:
    """Read RTF file"""
    try:
        import os
        import tempfile

        try:
            import pypandoc

            with tempfile.NamedTemporaryFile(delete=False, suffix=".rtf") as f:
                f.write(file.read())
                tmp = f.name

            try:
                return pypandoc.convert_file(tmp, "plain", format="rtf")
            finally:
                os.unlink(tmp)
        except ImportError:
            from striprtf.striprtf import rtf_to_text

            file.seek(0)
            return rtf_to_text(file.read().decode("utf-8"))

    except Exception as e:
        return f"Error: {e}"


def clean_text(text: str) -> str:
    import re
    import unicodedata

    text = unicodedata.normalize("NFKC", text)
    text = re.sub(r"[.\s\u00A0\u2009\u200B\u2024\u2027]{3,}", " ", text)

    def is_safe(c):
        cat = unicodedata.category(c)
        return cat[0] in "LNZ" or c in "\n\t @#/\\_-+=*&%$:.?!"

    text = "".join(c if is_safe(c) else " " for c in text)
    text = re.sub(r"\n\s+", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"(\b.{1,4}?\b)(?:\s*\1){2,}", r"\1", text)
    text = re.sub(r"(.)\1{2,}", r"\1\1", text)

    return text.strip()


def extract_text_from_file(file) -> str:
    """Extract text from any supported document type"""
    # Supported file types mapping
    SUPPORTED_TYPES = {
        # Text formats
        "text/plain": read_txt,
        "text/markdown": read_markdown,
        "text/html": read_html,
        # Microsoft Office
        "application/pdf": read_pdf,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": read_docx,  # .docx
        "application/msword": read_docx,  # .doc
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": read_excel,  # .xlsx
        "application/vnd.ms-excel": read_excel,  # .xls
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": read_pptx,  # .pptx
        "application/vnd.ms-powerpoint": read_pptx,  # .ppt
        # Data formats
        "text/csv": read_csv,
        "application/json": read_json,
        # Other
        "application/rtf": read_rtf,
        "text/rtf": read_rtf,
    }
    file_type = file.type

    # Try by MIME type
    if file_type in SUPPORTED_TYPES:
        return clean_text(SUPPORTED_TYPES[file_type](file))

    # Try by file extension as fallback
    file_ext = os.path.splitext(file.name)[1].lower()
    ext_mapping = {
        ".txt": read_txt,
        ".md": read_markdown,
        ".pdf": read_pdf,
        ".docx": read_docx,
        ".doc": read_docx,
        ".xlsx": read_excel,
        ".xls": read_excel,
        ".csv": read_csv,
        ".pptx": read_pptx,
        ".ppt": read_pptx,
        ".json": read_json,
        ".html": read_html,
        ".htm": read_html,
        ".rtf": read_rtf,
    }

    if file_ext in ext_mapping:
        return clean_text(ext_mapping[file_ext](file))
    return ""
